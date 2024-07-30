import pickle
import math
import os
import warnings
warnings.simplefilter(action='ignore', category=FutureWarning)
import pandas as pd
import docx
import pptx
import fitz 
from transformers import pipeline

print("\n🔃 Refreshing the Search Engine...\n")
import nltk
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')
try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords')
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from dotenv import load_dotenv
load_dotenv()

model_dir = os.getenv("MODEL_DIR")
qa_model_name = os.getenv("QA_MODEL_NAME")
STOPWORDS = set(stopwords.words("english"))

if not os.path.exists(model_dir):
    print(f"\n❌ Directory {model_dir} does not exist. Please set the correct MODEL_DIR in .env file.\n")
    exit()

qa_model = pipeline('question-answering', model=qa_model_name, tokenizer=qa_model_name)

# Global variables to hold the loaded model data
document_filenames = []
vocabulary = {}
postings = {}
document_frequency = {}
length = []
N = 0

# Function to load the model
def load_model(user="User", model_name="latest.pkl"):
    """Loads the Vector Space Model from the disk."""
    try:
        global document_filenames, vocabulary, postings, document_frequency, length, N
        with open(f'{model_dir}/{user}/{model_name}', 'rb') as f:
            model_data = pickle.load(f)
        document_filenames = model_data["document_filenames"]
        vocabulary = model_data["vocabulary"]
        postings = model_data["postings"]
        document_frequency = model_data["document_frequency"]
        length = model_data["length"]
        N = model_data["N"]
    except Exception as e:
        print(f"\n📂 Model not found in {model_dir}. Please build the model first using `buildModel.py`\n")
        print(f"Error: {e}")
        # exit()

def tokenize(document):
    """Returns a list whose elements are the separate terms in document."""
    terms = word_tokenize(document)
    terms = [term.lower() for term in terms if term not in STOPWORDS]
    return terms

def similarity(query, id):
    """Returns the cosine similarity between query and document id.
    Note that we don't bother dividing by the length of the query vector,
    since this doesn't make any difference to the ordering of search results.
    """
    similarity = 0.0  # Initialize similarity to 0
    for term in query:
        if term in vocabulary:
            # For every term in query which is also in vocabulary, calculate tf-idf score of the term and add to similarity
            similarity += term_frequency(term, id) * inverse_document_frequency(term)
    
    if length[id] == 0:
        return 0.0  # Return 0 if length of document vector is zero

    similarity = similarity / length[id]
    return similarity

def term_frequency(term, id):
    """Returns the term frequency of term in document id.
    If the term isn't in the document, then return 0.
    """
    if term in postings and id in postings[term]:
        return postings[term][id]
    else:
        return 0.0

def inverse_document_frequency(term):
    """Returns the inverse document frequency of term. 
    Note that if term isn't in the vocabulary then it returns 0, by convention.
    """
    if term in vocabulary:
        return math.log(N / document_frequency[term], 2)
    else:
        return 0.0

def read_document(filename):
    """Reads the content of a document and returns it as a string."""
    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext == '.txt':
            with open(filename, 'r', encoding='utf-8') as f:
                return f.read()
        elif ext == '.xlsx':
            df = pd.read_excel(filename)
            header = ' '.join(df.columns.values)
            values = ' '.join(df.astype(str).stack().tolist())
            return header + " " + values
        elif ext == '.docx':
            doc = docx.Document(filename)
            return ' '.join([para.text for para in doc.paragraphs])
        elif ext == '.pptx':
            prs = pptx.Presentation(filename)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return ' '.join(text)
        elif ext == '.pdf':
            doc = fitz.open(filename)
            text = []
            for page in doc:
                text.append(page.get_text())
            return ' '.join(text)
        else:
            return "Unsupported file format."
    except Exception as e:
        return f"Error reading document: {e}"

def get_snippet(document_filename, query, size=250):
    """Returns a snippet of text from the document where the query matches."""
    try:
        document = read_document(document_filename)
        terms = word_tokenize(document)
        terms_lower = [term.lower() for term in terms]

        for i in range(len(terms_lower)):
            if any(query_term in terms_lower[i] for query_term in query):
                start = max(0, i - size)
                end = min(len(terms), i + size)
                snippet = " ".join(terms[start:end])
                return snippet
        return "No snippet found."
    except Exception as e:
        return f"Error getting snippet: {e}"

def get_context(query, snippet):
    """Returns the context of the snippet by asking a question to the model."""

    snippet += " " + "Give me detailed answer to this question"
    QA_input = {
        'question': query,
        'context': snippet
    }
    res = qa_model(QA_input)
    try:
        answer = res['answer']
    except:
        answer = "No context found!"
    return answer
    
def print_scores(scores, query, original_query):
    """Prints scores in a tabular format with two columns (only used when __main__ is used)"""
    ct = len(scores)
    for (id, score) in scores:
        if score != 0.0:
            ct-=1
            snippet = get_snippet(document_filenames[id], query, 10)
            context = get_context(original_query, get_snippet(document_filenames[id], query, 250))
            print(f"🔢 Score: {str(score)[:5]}")
            print(f"📄 Document: {document_filenames[id]}")
            print(f"🔍 Quick: {context}")
            print(f"📚 Snippet: {snippet}\n")
    if ct == len(scores):
        print("🟪 No search results found! Check for possible typos in your query.\n")
    print("-" * 155, end="\n\n")

def do_search(query):
    """Asks the user what they would like to search for,
    and returns a list of relevant documents, in decreasing order of cosine similarity."""
    query = tokenize(str(query))

    scores = sorted(
        [(id, similarity(query, id)) for id in range(N)],
        key=lambda x: x[1],
        reverse=True,
    )
    print(f"⚡ Search job compeleted!\nScores: {scores}\n")
    return scores, query

def query_search():
    """Handles the search functionality."""
    while True:
        query = input("💻 Enter your search query: ")
        # Exit if query is empty, may be used if __main__ is used
        if query == "":
            print("\n👋 Goodbye! Exiting the engine...\n")
            exit()
        scores, query_terms = do_search(query)
        print()
        print_scores(scores, query_terms, original_query=query)

if __name__ == "__main__":
    """Main function to run the search engine."""
    load_model(user="admin", model_name="latest.pkl")
    query_search()