import math
import re
import os
import pickle
import time
from collections import defaultdict
from functools import reduce
from PyPDF2 import PdfReader
import docx
import pptx
import pandas as pd
from src.logger import *
import nltk
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')
try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords')
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from dotenv import load_dotenv
load_dotenv()

STOPWORDS = set(stopwords.words("english"))

# Global variables
document_filenames = dict()
N = 0
vocabulary = set()
postings = defaultdict(dict)
document_frequency = defaultdict(int)
length = defaultdict(float)
COUNT = 0

def get_all_files(directory):
    """Recursively gets all files from the given directory and subdirectories."""
    files = []
    for root, dirs, filenames in os.walk(directory):
        for filename in filenames:
            if filename.endswith(('.txt', '.pdf', '.docx', '.pptx', '.xlsx')):
                files.append(os.path.join(root, filename))
            else:
                log(f"🟨 Ignoring file: {filename} as it is not supported")
    return files

def read_file(filepath):
    """Reads the content of the file based on its extension."""
    global COUNT
    try:
        COUNT += 1
        if filepath.endswith('.txt'):
            with open(filepath, 'r') as file:
                return file.read()
        elif filepath.endswith('.pdf'):
            reader = PdfReader(filepath)
            text = ''
            for page in reader.pages:
                text += page.extract_text()
            return text
        elif filepath.endswith('.docx'):
            doc = docx.Document(filepath)
            return '\n'.join([para.text for para in doc.paragraphs])
        elif filepath.endswith('.pptx'):
            presentation = pptx.Presentation(filepath)
            text = ''
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + '\n'
            return text
        elif filepath.endswith('.xlsx'):
            df = pd.read_excel(filepath)
            return df.to_string()
        return ""
    except Exception as e:
        log(f"🟥 Error reading file: {filepath} - {e}")
        return ""

def initialize(directory, model_name, user="User"):
    """
    Initializes the vector space model by reading in a set of documents from a corpus (only when __main__ is used).
    """
    global CORPUS, text_dir, model_dir
    text_dir = directory
    model_dir = os.getenv("MODEL_DIR")
    if not os.path.exists(model_dir):
        os.makedirs(model_dir)

    log(f"🚀 Building in progress...")
    print("🚀 Building in progress...\n")
    clear_logs()
    get_corpus()
    initialize_terms_and_postings()
    initialize_document_frequencies()
    initialize_lengths()
    save_model(user, model_name)
    log(f"✅ Model built successfully and saved at {model_dir}/{user}/{model_name}")
    print(f"✅ Model built successfully and saved at {model_dir}/{user}/{model_name}\n")

def get_corpus():
    """Reads all the files in the corpus and stores them in the global variable document_filenames."""
    global document_filenames, N
    documents = get_all_files(text_dir)
    N = len(documents)
    # Dictionary having doc id as key and document name as value
    document_filenames = dict(zip(range(N), documents))

def initialize_terms_and_postings():
    """Reads in each document in document_filenames, splits it into a list of terms (i.e., tokenizes it),
    adds new terms to the global vocabulary, and adds the document to the posting list for each term,
    with value equal to the frequency of the term in the document.
    """
    global vocabulary, postings

    for id in document_filenames:
        document = read_file(document_filenames[id])
        document = remove_special_characters(document)
        document = remove_digits(document)
        terms = tokenize(document)
        unique_terms = set(terms)
        vocabulary = vocabulary.union(unique_terms) # Add unique terms to the vocabulary

        for term in unique_terms:
            postings[term][id] = terms.count(term)

def tokenize(document):
    """Returns a list whose elements are the separate terms in document."""
    terms = word_tokenize(document)
    terms = [term.lower() for term in terms if term not in STOPWORDS]
    return terms

def initialize_document_frequencies():
    """For each term in the vocabulary, count the number of documents it appears in,
    and store the value in document_frequency[term]
    """
    global document_frequency
    for term in vocabulary:
        document_frequency[term] = len(postings[term])

def initialize_lengths():
    """Computes the length for each document."""
    global length
    for id in document_filenames:
        l = 0
        for term in vocabulary:
            l += term_frequency(term, id) ** 2
        length[id] = math.sqrt(l)

def term_frequency(term, id):
    """Returns the term frequency of term in document id.
    If the term isn't in the document, then return 0.
    """
    if id in postings[term]:
        return postings[term][id]
    else:
        return 0.0

def intersection(sets):
    """Returns the intersection of all sets in the list sets.
    Requires that the list sets contains at least one element, otherwise it raises an error.
    """
    return reduce(set.intersection, [s for s in sets])

def remove_special_characters(text):
    """ Removes special characters using regex substitution """
    regex = re.compile(r"[^a-zA-Z0-9\s]")
    return re.sub(regex, "", text)

def remove_digits(text):
    """ Removes digits using regex substitution """
    regex = re.compile(r"\d")
    return re.sub(regex, "", text)

def save_model(user, model_name):
    """Saves the model data to a pickle file."""
    model_data = {
        "document_filenames": document_filenames,
        "vocabulary": vocabulary,
        "postings": postings,
        "document_frequency": document_frequency,
        "length": length,
        "N": N
    }
    os.makedirs(f"{model_dir}/{user}", exist_ok=True)
    with open(f'{model_dir}{user}/{model_name}', 'wb') as f:
        pickle.dump(model_data, f)

if __name__ == "__main__":
    """Main function to build the vector space model."""
    directory = input("\n📂 Enter the directory path for which you want to build the model: ")
    model_name = input("\n📄 Enter the model name to save the model as (default: latest.pkl): ")
    user = input("\n👤 Enter the username to save the model as (default: self): ")
    if model_name == "":
        model_name = "latest.pkl"
    if user == "":
        user = "self"
    print("\n🚀 Building in progress...\n")
    stime = time.time()
    initialize(directory, user, model_name)
    print(f"\n⏳ Time taken to build the model: {round(time.time() - stime, 2)} seconds. Files read: {COUNT}\n")